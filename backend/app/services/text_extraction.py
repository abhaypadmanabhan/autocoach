"""Text extraction services for PDF and PPTX files."""

import logging
from io import BytesIO

from pypdf import PdfReader
from pptx import Presentation

logger = logging.getLogger(__name__)


def sanitize_text(text: str) -> str:
    """
    Sanitize text by removing null characters and problematic unicode.

    Args:
        text: The raw text to sanitize.

    Returns:
        Cleaned text with problematic characters removed.
    """
    if not text:
        return text

    # Remove null characters
    text = text.replace('\x00', '')
    text = text.replace('\u0000', '')

    # Remove other control characters except newlines (\n), tabs (\t), and carriage returns (\r)
    # Control characters are in range 0x00-0x1F and 0x7F
    cleaned_chars = []
    for char in text:
        code = ord(char)
        # Allow printable characters (32-126), newlines (10), tabs (9), carriage returns (13)
        if (32 <= code <= 126) or code in (9, 10, 13):
            cleaned_chars.append(char)
        # Skip other control characters (including null 0x00)

    return ''.join(cleaned_chars)


def extract_text_from_pdf(file_bytes: bytes) -> tuple[str, int]:
    """
    Extract text from a PDF file.

    Args:
        file_bytes: The raw bytes of the PDF file.

    Returns:
        A tuple of (full_text, page_count). Returns empty string and 0 if extraction fails.
    """
    try:
        pdf_file = BytesIO(file_bytes)
        reader = PdfReader(pdf_file)

        pages_text = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages_text.append(text)

        full_text = "\n".join(pages_text)
        full_text = sanitize_text(full_text)
        page_count = len(reader.pages)

        logger.info(f"Extracted {page_count} pages from PDF, {len(full_text)} characters")
        return full_text, page_count
    except Exception as e:
        logger.error(f"Failed to extract text from PDF: {e}")
        return "", 0


def extract_text_from_pptx(file_bytes: bytes) -> tuple[str, int]:
    """
    Extract text from a PowerPoint file.

    Args:
        file_bytes: The raw bytes of the PPTX file.

    Returns:
        A tuple of (full_text, slide_count). Returns empty string and 0 if extraction fails.
    """
    try:
        pptx_file = BytesIO(file_bytes)
        presentation = Presentation(pptx_file)

        slides_text = []
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        paragraph_text = paragraph.text.strip()
                        if paragraph_text:
                            slide_texts.append(paragraph_text)
            if slide_texts:
                slides_text.append(f"\n".join(slide_texts))

        full_text = "\n\n".join(slides_text)
        full_text = sanitize_text(full_text)
        slide_count = len(presentation.slides)

        logger.info(f"Extracted {slide_count} slides from PPTX, {len(full_text)} characters")
        return full_text, slide_count
    except Exception as e:
        logger.error(f"Failed to extract text from PPTX: {e}")
        return "", 0
